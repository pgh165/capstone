"""
AI 기반 학습 피로 관리 시스템 — PowerPoint 생성 스크립트
실행: python make_ppt.py  →  presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree

# ── 애니메이션 헬퍼 ──────────────────────────────────────────

def add_transition(slide, style="fade", dur_ms=600):
    """슬라이드 전환 효과 추가. spd/dur 중 dur만 사용(상호배타)."""
    sld = slide._element
    trans = etree.SubElement(sld, qn("p:transition"))
    trans.set("dur", str(dur_ms))   # spd 와 동시 사용 금지
    if style == "fade":
        etree.SubElement(trans, qn("p:fade"))
    elif style == "push":
        push = etree.SubElement(trans, qn("p:push"))
        push.set("dir", "l")
    elif style == "wipe":
        wipe = etree.SubElement(trans, qn("p:wipe"))
        wipe.set("dir", "l")




# ── 색상 팔레트 ──────────────────────────────────────────────
BG_DARK    = RGBColor(0x0A, 0x0C, 0x14)   # 슬라이드 배경
SURFACE    = RGBColor(0x12, 0x15, 0x1F)
CARD       = RGBColor(0x1A, 0x1E, 0x2E)
ACCENT     = RGBColor(0x5C, 0x8D, 0xFF)   # 파란 강조
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
YELLOW     = RGBColor(0xFF, 0xC1, 0x07)
ORANGE     = RGBColor(0xFF, 0x70, 0x43)
RED        = RGBColor(0xF4, 0x43, 0x36)
PURPLE     = RGBColor(0xAB, 0x47, 0xBC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x8A, 0x8F, 0xAA)
TEXT       = RGBColor(0xE8, 0xEA, 0xF4)

# 슬라이드 크기: 16:9 와이드
W = Inches(13.33)
H = Inches(7.5)


def prs_init():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── 헬퍼 ─────────────────────────────────────────────────────
def fill_bg(slide, color=BG_DARK):
    """슬라이드 배경색 채우기."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, radius=False):
    """단색 사각형 추가."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return txBox


def add_title_box(slide, title, subtitle=None):
    """상단 제목 영역."""
    # 배경 라인
    bar = add_rect(slide, Inches(0), Inches(0), W, Pt(4), ACCENT)
    add_textbox(slide, title,
                Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(0.82), Inches(12), Inches(0.4),
                    font_size=14, color=MUTED, align=PP_ALIGN.LEFT)


def add_card(slide, left, top, width, height, title=None, body_lines=None,
             title_color=ACCENT, body_color=TEXT, icon=None):
    """카드 박스."""
    add_rect(slide, left, top, width, height, CARD)
    # 테두리 효과 (상단 라인)
    add_rect(slide, left, top, width, Pt(3), title_color)

    y = top + Inches(0.15)
    if icon and title:
        add_textbox(slide, f"{icon}  {title}",
                    left + Inches(0.15), y, width - Inches(0.3), Inches(0.4),
                    font_size=14, bold=True, color=title_color)
    elif title:
        add_textbox(slide, title,
                    left + Inches(0.15), y, width - Inches(0.3), Inches(0.4),
                    font_size=14, bold=True, color=title_color)
    if body_lines:
        body_top = y + Inches(0.45)
        for i, line in enumerate(body_lines):
            add_textbox(slide, line,
                        left + Inches(0.2), body_top + i * Inches(0.32),
                        width - Inches(0.4), Inches(0.35),
                        font_size=11, color=body_color)


def page_num(slide, current, total):
    add_textbox(slide, f"{current} / {total}",
                Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
                font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ── 슬라이드 생성 ─────────────────────────────────────────────
TOTAL = 10

def slide01_title(prs):
    """표지."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_bg(sld)

    # 그라데이션 느낌 — 좌측 세로 바
    add_rect(sld, Inches(0), Inches(0), Inches(0.08), H, ACCENT)

    # 중앙 상단 태그
    add_textbox(sld, "Capstone Project",
                Inches(1.2), Inches(1.2), Inches(11), Inches(0.5),
                font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # 메인 타이틀
    add_textbox(sld,
                "AI 개인 맞춤형 포모도로 타이머",
                Inches(1.2), Inches(1.9), Inches(11), Inches(1.0),
                font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 부제목
    add_textbox(sld,
                "공부 중 졸음·집중력 저하 실시간 감지 및 맞춤형 휴식 관리 시스템",
                Inches(1.2), Inches(3.0), Inches(11), Inches(0.6),
                font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # 구분선
    add_rect(sld, Inches(4.5), Inches(3.75), Inches(4.3), Pt(2), MUTED)

    # 하단 정보
    add_textbox(sld, "MediaPipe  ·  Ollama LLM  ·  MeloTTS  ·  Django  ·  OpenCV",
                Inches(1.2), Inches(4.1), Inches(11), Inches(0.4),
                font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

    add_textbox(sld, "2026",
                Inches(1.2), Inches(6.8), Inches(11), Inches(0.4),
                font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

    page_num(sld, 1, TOTAL)
    # 전환: 페이드
    add_transition(sld, "fade", dur_ms=800)
    return sld


def slide02_problem(prs):
    """문제 정의."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "문제 정의", "공부할 때 피로를 제대로 관리하기 어려운 이유")

    stats = [
        ("😴", "피로 누적의 악순환",
         "장시간 학습 → 집중력 저하 → 효율 감소\n"
         "졸음이 와도 '조금만 더' 버티다 역효과\n"
         "피로도를 스스로 정확히 파악하기 어려움"),
        ("⏱", "기존 포모도로의 한계",
         "고정된 25분 작업 / 5분 휴식 패턴\n"
         "개인 컨디션·피로 상태 전혀 반영 안 됨\n"
         "컨디션 좋을 땐 너무 짧고, 나쁠 땐 너무 길다"),
        ("🤖", "AI 맞춤 관리의 필요성",
         "웹캠으로 실시간 졸음 수치 객관화\n"
         "피로 이력 학습 → 개인별 최적 휴식 시점\n"
         "LLM이 상황에 맞는 회복 방법 직접 조언"),
    ]

    for i, (icon, title, body) in enumerate(stats):
        x = Inches(0.4) + i * Inches(4.25)
        add_card(sld, x, Inches(1.45), Inches(4.0), Inches(5.6),
                 title=f"{icon}  {title}", body_lines=body.split("\n"),
                 title_color=ACCENT)

    page_num(sld, 2, TOTAL)
    add_transition(sld, "push")
    return sld


def slide03_architecture(prs):
    """시스템 아키텍처."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "시스템 아키텍처", "WSL2 직접 실행 구조 — Docker 없이 웹캠 직접 연결")

    # 파이프라인 박스들
    pipeline = [
        ("📷 Camera",   "USB 웹캠\nMJPEG 캡처",          GREEN),
        ("🔍 Detector", "MediaPipe\nFace Mesh",           ACCENT),
        ("📊 Judge",    "EAR / MAR\nHead Pose → Score",   YELLOW),
        ("🤖 AI Judge", "Ollama LLM\n보정 판정",           PURPLE),
        ("🔔 Alert",    "히스테리시스\n4단계 경보",         ORANGE),
    ]

    box_w = Inches(2.2)
    box_h = Inches(1.5)
    gap   = Inches(0.2)
    start_x = Inches(0.35)
    top_y = Inches(1.6)

    for i, (title, body, color) in enumerate(pipeline):
        x = start_x + i * (box_w + gap)
        add_card(sld, x, top_y, box_w, box_h,
                 title=title, body_lines=body.split("\n"),
                 title_color=color)
        # 화살표
        if i < len(pipeline) - 1:
            ax = x + box_w + Inches(0.02)
            add_textbox(sld, "→", ax, top_y + Inches(0.55), gap, Inches(0.4),
                        font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # 하단 두 박스
    side = [
        ("💤 Fatigue",   "피로도 누적 관리\n작업 시간 + 졸음 빈도",   Inches(0.35),  RED),
        ("🍅 Pomodoro",  "AI 동적 타이머\n피로·졸음 기반 인터벌 조정", Inches(3.15),  GREEN),
        ("💬 LLM Coach", "Ollama gemma4:26b\n개인화 코칭 메시지",      Inches(5.95),  ACCENT),
        ("🔊 Voice",     "MeloTTS (한국어)\n로컬 AI 음성 출력",        Inches(8.75),  PURPLE),
        ("📈 Dashboard", "Django + 실시간\n대시보드 웹 UI",            Inches(11.05), YELLOW),
    ]

    bot_y = Inches(3.5)
    for title, body, x, color in side:
        add_card(sld, x, bot_y, Inches(2.5), Inches(1.4),
                 title=title, body_lines=body.split("\n"),
                 title_color=color)

    # 연결선 설명
    add_textbox(sld, "↑ 결과 전달",
                Inches(0.35), Inches(3.2), Inches(10), Inches(0.3),
                font_size=10, color=MUTED)

    page_num(sld, 3, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide04_drowsiness(prs):
    """졸음 감지 알고리즘."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "졸음 감지 알고리즘", "EAR·MAR·Head Pose 다중 지표 융합 + 제곱 변환")

    metrics = [
        ("👁 EAR (눈 감김)", ACCENT, [
            "Eye Aspect Ratio = 눈 높이 / 눈 너비",
            "임계값 0.2 이하 → 눈 감김 판정",
            "2초 이상 지속 시 높은 점수 부여",
            "가중치 W1 = 0.45 (가장 직접적)",
        ]),
        ("😮 MAR (하품)", YELLOW, [
            "Mouth Aspect Ratio = 입 개방 비율",
            "임계값 0.6 이상 → 하품 판정",
            "3분 내 3회 이상 → 졸음 전조",
            "가중치 W2 = 0.30",
        ]),
        ("🤕 Head Pose (고개 숙임)", ORANGE, [
            "피치(Pitch) 값 기반 — 고개 앞으로 숙임",
            "요(Yaw) 제외 — 옆보기는 졸음 아님",
            "15° 데드존 → 자연 움직임 오감지 방지",
            "가중치 W3 = 0.25",
        ]),
    ]

    for i, (title, color, lines) in enumerate(metrics):
        x = Inches(0.35) + i * Inches(4.3)
        add_card(sld, x, Inches(1.45), Inches(4.1), Inches(3.5),
                 title=title, body_lines=lines, title_color=color)

    # 수식 박스
    formula_box = add_rect(sld, Inches(0.35), Inches(5.3), Inches(12.63), Inches(1.8), CARD)
    add_rect(sld, Inches(0.35), Inches(5.3), Inches(12.63), Pt(3), PURPLE)

    add_textbox(sld, "최종 졸음 점수 산출",
                Inches(0.55), Inches(5.35), Inches(5), Inches(0.4),
                font_size=13, bold=True, color=PURPLE)

    add_textbox(sld,
                "raw = 0.45×EAR점수 + 0.30×MAR점수 + 0.25×Head점수",
                Inches(0.55), Inches(5.8), Inches(8), Inches(0.4),
                font_size=13, color=TEXT)

    add_textbox(sld,
                "score = raw² / 100   (제곱 변환: 낮은 수치 억제, 높은 수치 민감)",
                Inches(0.55), Inches(6.15), Inches(10), Inches(0.4),
                font_size=13, color=ACCENT)

    add_textbox(sld,
                "→ EMA(α=0.15) 스무딩 적용 → 히스테리시스 경보 판정",
                Inches(0.55), Inches(6.5), Inches(10), Inches(0.4),
                font_size=12, color=MUTED)

    page_num(sld, 4, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide05_alert(prs):
    """경보 시스템."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "4단계 경보 시스템", "히스테리시스 적용 — 경보 진동(oscillation) 방지")

    levels = [
        ("L0  정상",   "0 ~ 40",  GREEN,  "얼굴 정상, 집중 작업 중\n경보 없음"),
        ("L1  주의",   "41 ~ 70", YELLOW, "가벼운 졸음 감지\n'주의하세요' TTS 발화"),
        ("L2  경고",   "71 ~ 85", ORANGE, "졸음 수준 높음\n'경고! 잠시 쉬어가세요'"),
        ("L3  위험",   "86 ~ 100",RED,    "즉각적 휴식 필요\n포모도로 강제 휴식 전환"),
    ]

    for i, (label, score_range, color, desc) in enumerate(levels):
        x = Inches(0.35) + i * Inches(3.22)
        add_card(sld, x, Inches(1.45), Inches(3.05), Inches(2.5),
                 title=label, body_lines=[f"점수: {score_range}", *desc.split("\n")],
                 title_color=color)

    # 히스테리시스 설명
    add_rect(sld, Inches(0.35), Inches(4.3), Inches(12.63), Inches(1.3), CARD)
    add_rect(sld, Inches(0.35), Inches(4.3), Inches(12.63), Pt(3), ACCENT)
    add_textbox(sld, "히스테리시스 (Hysteresis) — 경보 진동 방지",
                Inches(0.55), Inches(4.35), Inches(8), Inches(0.4),
                font_size=13, bold=True, color=ACCENT)
    add_textbox(sld,
                "• 상승: 규칙 기반 임계값 초과 시 즉시 레벨 올림\n"
                "• 하강: 더 낮은 임계값(L1→0: 30점, L2→1: 58점, L3→2: 75점) 통과해야 내림\n"
                "• 효과: 점수가 경계값 근처에서 오르락내리락해도 경보가 반복 발화되지 않음",
                Inches(0.55), Inches(4.78), Inches(12.3), Inches(0.8),
                font_size=11, color=TEXT)

    # AI Judge 설명
    add_rect(sld, Inches(0.35), Inches(5.75), Inches(12.63), Inches(1.3), CARD)
    add_rect(sld, Inches(0.35), Inches(5.75), Inches(12.63), Pt(3), PURPLE)
    add_textbox(sld, "AI Judge (Ollama gemma4:e4b) — 규칙 기반 보정",
                Inches(0.55), Inches(5.8), Inches(8), Inches(0.4),
                font_size=13, bold=True, color=PURPLE)
    add_textbox(sld,
                "• 5초 주기로 EAR/MAR/Pitch 원시 수치를 LLM에 전달 → 졸음 점수 재판정\n"
                "• 규칙 기반 점수와 30점 이상 차이 시 중간값으로 완화 (급격한 변동 방지)\n"
                "• 경량 모델(4B) 사용 — 평균 응답 3~5초, 메인 루프 비차단",
                Inches(0.55), Inches(6.23), Inches(12.3), Inches(0.8),
                font_size=11, color=TEXT)

    page_num(sld, 5, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide06_pomodoro(prs):
    """AI 포모도로 타이머."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "AI 개인 맞춤형 포모도로 타이머", "피로도·졸음 점수 기반 동적 작업/휴식 인터벌 조정")

    # 왼쪽: 동작 원리
    add_card(sld, Inches(0.35), Inches(1.45), Inches(6.0), Inches(5.6),
             title="⚙️  동적 인터벌 조정 원리", title_color=GREEN,
             body_lines=[
                 "기본 작업 인터벌: 25분",
                 "",
                 "컨디션 좋음 (피로↓·졸음↓):",
                 "  → 인터벌 연장 (최대 40분)",
                 "",
                 "컨디션 나쁨 (피로↑·졸음↑):",
                 "  → 인터벌 단축 (최소 10분)",
                 "",
                 "위험 수준 감지:",
                 "  → 즉시 강제 휴식 전환",
                 "",
                 "4사이클 완료:",
                 "  → 긴 휴식 15분 자동 적용",
             ])

    # 오른쪽: 개인화 회복 프로필
    add_card(sld, Inches(6.65), Inches(1.45), Inches(6.33), Inches(2.6),
             title="🧬  개인화 회복 프로필 학습", title_color=PURPLE,
             body_lines=[
                 "DB의 recovery_actions 이력 분석",
                 "가이드별 성공률 계산 (눈 휴식, 스트레칭, 걷기...)",
                 "개인에게 효과적인 가이드 우선 추천",
                 "최소 3건 이상 이력 축적 시 개인화 적용",
             ])

    add_card(sld, Inches(6.65), Inches(4.3), Inches(6.33), Inches(2.75),
             title="🔔  TTS 음성 안내 (MeloTTS)", title_color=ACCENT,
             body_lines=[
                 "포모도로 시작: '25분 집중해볼까요?'",
                 "휴식 전환:    '휴식 시간입니다. 5분 쉬어갑니다.'",
                 "복귀:         '휴식 끝! 다음 작업은 N분입니다.'",
                 "",
                 "→ 로컬 AI TTS — 인터넷 연결 불필요",
                 "→ 비차단 큐 기반 — 메인 루프 영향 없음",
             ])

    page_num(sld, 6, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide07_llm_coach(prs):
    """LLM 코치."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "LLM 개인화 코칭", "Ollama gemma4:26b — 로컬 AI 맞춤 조언 생성")

    # 흐름
    flow = [
        ("🚨 트리거",    GREEN,  "포모도로 휴식 전환 시\n경고·위험 단계 감지 시"),
        ("📋 컨텍스트",  ACCENT, "피로 단계·점수\n연속 작업 시간\n개인 회복 이력"),
        ("🤖 LLM 처리", PURPLE, "gemma4:26b\nThinking 모드\n1,500 토큰"),
        ("💬 출력",      YELLOW, "3~4문장 대화체\n공감 → 조언 → 응원\n마크다운 금지"),
        ("🔊 TTS 발화",  ORANGE, "MeloTTS 한국어\n음성 출력"),
    ]

    box_w = Inches(2.3)
    for i, (title, color, body) in enumerate(flow):
        x = Inches(0.35) + i * (box_w + Inches(0.18))
        add_card(sld, x, Inches(1.45), box_w, Inches(2.5),
                 title=title, body_lines=body.split("\n"), title_color=color)
        if i < len(flow) - 1:
            ax = x + box_w + Inches(0.01)
            add_textbox(sld, "→", ax, Inches(2.3), Inches(0.17), Inches(0.4),
                        font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 예시 메시지 박스
    add_rect(sld, Inches(0.35), Inches(4.2), Inches(12.63), Inches(2.85), CARD)
    add_rect(sld, Inches(0.35), Inches(4.2), Inches(12.63), Pt(3), PURPLE)
    add_textbox(sld, "💬  실제 출력 예시 (위험 단계)",
                Inches(0.55), Inches(4.25), Inches(8), Inches(0.4),
                font_size=13, bold=True, color=PURPLE)
    add_textbox(sld,
                '"님, 계속되는 졸음 때문에 정말 많이 피곤하시겠어요.\n'
                '지금 바로 눈을 감고 1분만 눈 주변을 부드럽게 마사지해 보세요.\n'
                '잠시 쉬고 나면 훨씬 개운해질 테니 조금만 더 힘내세요."',
                Inches(0.55), Inches(4.72), Inches(12.3), Inches(1.3),
                font_size=13, color=TEXT, italic=True)

    add_textbox(sld,
                "• 비동기 처리 — 응답 대기 중에도 메인 루프 정상 작동\n"
                "• 300초 쿨다운 — 과도한 LLM 호출 방지",
                Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.55),
                font_size=11, color=MUTED)

    page_num(sld, 7, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide08_dashboard(prs):
    """대시보드."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "실시간 모니터링 대시보드", "Django 기반 웹 UI — 브라우저에서 바로 확인")

    pages = [
        ("📊  실시간 모니터링", ACCENT, [
            "1초 폴링 — 졸음·피로·EAR·MAR 수치 실시간 갱신",
            "Chart.js 라인 차트 — 최근 100포인트 히스토리",
            "sessionStorage — 페이지 이동 후에도 데이터 유지",
            "브라우저 Push 알림 (Notification API)",
            "포모도로 상태·완료 횟수 표시",
            "프로그램 종료 시 '연결 끊김' 오버레이",
        ]),
        ("📈  일별 리포트", YELLOW, [
            "날짜 선택기 — 과거 데이터 조회",
            "페이지네이션 — 10건씩 로그 분할",
            "피로도·졸음 집계 차트",
            "DB 저장 주기: 5초마다 자동 저장",
        ]),
        ("⚙️  설정 페이지", GREEN, [
            "EAR·MAR 임계값 실시간 조정",
            "졸음 점수 가중치(W1/W2/W3) 변경",
            "AJAX 저장 — 재시작 불필요",
        ]),
    ]

    for i, (title, color, lines) in enumerate(pages):
        x = Inches(0.35) + i * Inches(4.3)
        add_card(sld, x, Inches(1.45), Inches(4.1), Inches(5.6),
                 title=title, body_lines=lines, title_color=color)

    page_num(sld, 8, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide09_tech_stack(prs):
    """기술 스택."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "기술 스택 요약", "전 구성 요소 로컬 실행 — 클라우드·인터넷 연결 불필요")

    stacks = [
        ("👁  감지 레이어", ACCENT, [
            "OpenCV — 웹캠 MJPEG 캡처",
            "MediaPipe Face Mesh — 468개 랜드마크",
            "EAR / MAR / Head Pose 실시간 계산",
        ]),
        ("🤖  AI 레이어", PURPLE, [
            "Ollama — 로컬 LLM 런타임",
            "gemma4:e4b (9GB) — 빠른 AI 판정",
            "gemma4:26b (17GB) — 고품질 코칭",
        ]),
        ("🔊  음성 레이어", GREEN, [
            "MeloTTS — 로컬 한국어 AI TTS",
            "비차단 큐 — 메인 루프 영향 없음",
            "PowerShell 브릿지 (WSL2 오디오)",
        ]),
        ("🌐  웹 레이어", YELLOW, [
            "Django 4.x — 웹 프레임워크",
            "MariaDB (Docker) — 데이터 저장",
            "Chart.js + sessionStorage",
        ]),
        ("⚡  실행 환경", ORANGE, [
            "WSL2 Ubuntu 24.04",
            "Python 3.11 venv",
            "Docker Compose (DB only)",
        ]),
    ]

    for i, (title, color, lines) in enumerate(stacks):
        row, col = divmod(i, 3)
        x = Inches(0.35) + col * Inches(4.3)
        y = Inches(1.45) + row * Inches(2.7)
        add_card(sld, x, y, Inches(4.1), Inches(2.4),
                 title=title, body_lines=lines, title_color=color)

    page_num(sld, 9, TOTAL)
    add_transition(sld, 'push')
    return sld


def slide10_demo(prs):
    """시연 & 결론."""
    sld = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(sld)
    add_title_box(sld, "시연 결과 & 결론", "demo.py — 카메라 없이 전체 파이프라인 단계별 시연")

    # 시연 단계
    stages = [
        ("1 정상",  GREEN,  "졸음 < 1\n포모도로 시작"),
        ("2 주의",  YELLOW, "졸음 41.6\n'주의' TTS"),
        ("3 경고",  ORANGE, "졸음 77.4\n'경고' TTS"),
        ("4 위험",  RED,    "졸음 92.2\n'위험' TTS\nLLM 요청"),
        ("5 휴식",  PURPLE, "포모도로\n휴식 전환\n코칭 출력"),
        ("6 회복",  GREEN,  "졸음 < 15\n정상 복귀"),
    ]

    for i, (label, color, body) in enumerate(stages):
        x = Inches(0.35) + i * Inches(2.15)
        add_card(sld, x, Inches(1.45), Inches(2.0), Inches(2.4),
                 title=label, body_lines=body.split("\n"), title_color=color)
        if i < len(stages) - 1:
            add_textbox(sld, "→",
                        x + Inches(2.0), Inches(2.3), Inches(0.15), Inches(0.4),
                        font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 성과
    add_rect(sld, Inches(0.35), Inches(4.1), Inches(12.63), Inches(1.6), CARD)
    add_rect(sld, Inches(0.35), Inches(4.1), Inches(12.63), Pt(3), GREEN)
    add_textbox(sld, "✅  주요 성과",
                Inches(0.55), Inches(4.15), Inches(8), Inches(0.4),
                font_size=13, bold=True, color=GREEN)
    add_textbox(sld,
                "• 공부 중 졸음·피로를 수치로 객관화 — 본인 컨디션을 실시간으로 파악 가능\n"
                "• AI 포모도로로 컨디션에 맞는 휴식 시점 자동 제안 — '버티기' 대신 '현명한 휴식'\n"
                "• LLM 코치가 상황별 회복 방법을 말로 직접 안내 (눈 휴식, 스트레칭 등)\n"
                "• 전 구성 요소 로컬 실행 — 인터넷 없이 프라이버시 보호",
                Inches(0.55), Inches(4.58), Inches(12.3), Inches(1.0),
                font_size=11, color=TEXT)

    # 개선 방향
    add_rect(sld, Inches(0.35), Inches(5.85), Inches(12.63), Inches(1.2), CARD)
    add_rect(sld, Inches(0.35), Inches(5.85), Inches(12.63), Pt(3), MUTED)
    add_textbox(sld, "🔮  향후 개선 방향",
                Inches(0.55), Inches(5.9), Inches(8), Inches(0.4),
                font_size=13, bold=True, color=MUTED)
    add_textbox(sld,
                "학습 집중도 점수 추가 (시선 추적)  ·  "
                "GPU 가속으로 AI 판정 지연 단축  ·  "
                "학습 이력 기반 장기 피로 패턴 분석",
                Inches(0.55), Inches(6.32), Inches(12.3), Inches(0.6),
                font_size=11, color=MUTED)

    page_num(sld, 10, TOTAL)
    add_transition(sld, 'push')
    return sld


# ── 메인 ─────────────────────────────────────────────────────
def main():
    prs = prs_init()

    slide01_title(prs)
    slide02_problem(prs)
    slide03_architecture(prs)
    slide04_drowsiness(prs)
    slide05_alert(prs)
    slide06_pomodoro(prs)
    slide07_llm_coach(prs)
    slide08_dashboard(prs)
    slide09_tech_stack(prs)
    slide10_demo(prs)

    out = "presentation.pptx"
    prs.save(out)
    print(f"✅  저장 완료: {out}  ({TOTAL}슬라이드)")


if __name__ == "__main__":
    main()
